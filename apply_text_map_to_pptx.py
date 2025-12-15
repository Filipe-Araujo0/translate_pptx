#!/usr/bin/env python3
"""Clone a PPTX while replacing every mapped text entry with translated content."""

from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any, Dict, List, Sequence

from pptx import Presentation  # type: ignore[import-not-found]
from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore[import-not-found]

Entry = Dict[str, Any]


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description=(
            "Apply a translated text-map JSON onto a PPTX, writing a new PPTX with "
            "all mapped strings replaced in-place."
        )
    )
    parser.add_argument(
        "--pptx",
        required=True,
        help="Source PPTX file that matches the text map.",
    )
    parser.add_argument(
        "--translated-map",
        required=True,
        help="JSON file produced from build_text_map/apply_translated_texts containing translated entries.",
    )
    parser.add_argument(
        "--output",
        required=True,
        help="Destination PPTX path for the translated copy.",
    )
    return parser.parse_args()


def load_json(path: Path) -> Dict[str, Any]:
    data = json.loads(path.read_text())
    assert isinstance(data, dict), "Map JSON must be an object"
    assert isinstance(data.get("entries"), list), "Map JSON missing entries array"
    return data


def resolve_shape(slide, chain: Sequence[int]):
    assert chain, "shape_index_chain must not be empty"
    shapes = slide.shapes
    target = None
    for depth, idx in enumerate(chain, start=1):
        assert 1 <= idx <= len(shapes), (
            f"Shape index {idx} out of range at depth {depth}"
        )
        target = shapes[idx - 1]
        if depth < len(chain):
            assert target.shape_type == MSO_SHAPE_TYPE.GROUP, (
                "Intermediate shape must be a group"
            )
            shapes = target.shapes
    assert target is not None, "Failed to resolve shape"
    return target


def set_run_text(text_frame, paragraph_index: int, run_index: int, text: str) -> None:
    paragraphs = text_frame.paragraphs
    assert 1 <= paragraph_index <= len(paragraphs), "paragraph_index out of range"
    paragraph = paragraphs[paragraph_index - 1]
    runs = paragraph.runs
    assert 1 <= run_index <= len(runs), "run_index out of range"
    runs[run_index - 1].text = text


def apply_text_frame_entry(shape, entry: Entry) -> None:
    assert getattr(shape, "has_text_frame", False), "Shape lacks text_frame"
    set_run_text(
        shape.text_frame, entry["paragraph_index"], entry["run_index"], entry["text"]
    )


def apply_table_cell_entry(shape, entry: Entry) -> None:
    assert getattr(shape, "has_table", False), "Shape lacks table"
    row_idx = entry.get("table_row")
    col_idx = entry.get("table_col")
    assert isinstance(row_idx, int) and isinstance(col_idx, int), (
        "Missing table indices"
    )
    table = shape.table
    assert 1 <= row_idx <= len(table.rows), "table_row out of range"
    assert 1 <= col_idx <= len(table.columns), "table_col out of range"
    cell = table.cell(row_idx - 1, col_idx - 1)
    set_run_text(
        cell.text_frame, entry["paragraph_index"], entry["run_index"], entry["text"]
    )


def get_chart_nodes(chart_part, cache: Dict[str, List[Any]]) -> List[Any]:
    key = chart_part.partname
    if key not in cache:
        cache[key] = chart_part._element.xpath(".//a:t")
    return cache[key]


def apply_chart_entry(shape, entry: Entry, cache: Dict[str, List[Any]]) -> None:
    assert getattr(shape, "has_chart", False), "Shape lacks chart"
    chart_part = shape.chart.part
    assert chart_part.partname == entry["chart_partname"], "Chart part mismatch"
    nodes = get_chart_nodes(chart_part, cache)
    chart_idx = entry["chart_text_index"]
    assert isinstance(chart_idx, int) and chart_idx >= 1, "Invalid chart_text_index"
    assert chart_idx <= len(nodes), "chart_text_index out of range"
    nodes[chart_idx - 1].text = entry["text"]


def apply_entry(presentation, entry: Entry, chart_cache: Dict[str, List[Any]]) -> None:
    slide_index = entry["slide_index"]
    assert isinstance(slide_index, int) and 1 <= slide_index <= len(
        presentation.slides
    ), "slide_index out of range"
    slide = presentation.slides[slide_index - 1]
    shape_chain = entry["shape_index_chain"]
    assert isinstance(shape_chain, list), "shape_index_chain missing"
    shape = resolve_shape(slide, shape_chain)
    container = entry["container"]
    assert isinstance(container, str), "container must be string"
    if container == "text_frame":
        apply_text_frame_entry(shape, entry)
    elif container == "table_cell":
        apply_table_cell_entry(shape, entry)
    elif container == "chart_part":
        apply_chart_entry(shape, entry, chart_cache)
    else:
        raise ValueError(f"Unsupported container type: {container}")


def apply_text_map(
    pptx_path: Path, translated_map_path: Path, output_path: Path
) -> None:
    presentation = Presentation(pptx_path)
    text_map = load_json(translated_map_path)
    entries = text_map["entries"]
    expected_entry_count = text_map.get("entry_count")
    if expected_entry_count is not None:
        assert isinstance(expected_entry_count, int), (
            "entry_count must be an integer when present"
        )
        assert expected_entry_count == len(entries), (
            "entry_count does not match actual number of entries"
        )
    chart_cache: Dict[str, List[Any]] = {}
    for entry in entries:
        assert isinstance(entry, dict), "Entry must be object"
        apply_entry(presentation, entry, chart_cache)
    presentation.save(output_path)


def main() -> None:
    args = parse_args()
    pptx_path = Path(args.pptx)
    translated_map_path = Path(args.translated_map)
    output_path = Path(args.output)
    assert pptx_path.exists(), f"PPTX not found: {pptx_path}"
    assert translated_map_path.exists(), (
        f"Translated map not found: {translated_map_path}"
    )
    apply_text_map(pptx_path, translated_map_path, output_path)


if __name__ == "__main__":
    main()
