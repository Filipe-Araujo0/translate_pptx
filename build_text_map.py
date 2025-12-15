#!/usr/bin/env python3
"""Generate a JSON map describing every text-bearing element in a PPTX file."""

from __future__ import annotations

import argparse
import json
from pathlib import Path
from typing import Any, Dict, Iterable, List, Sequence

from pptx import Presentation  # type: ignore[import-not-found]
from pptx.enum.shapes import MSO_SHAPE_TYPE  # type: ignore[import-not-found]


def parse_args() -> argparse.Namespace:
    parser = argparse.ArgumentParser(
        description=(
            "Extract every text snippet (shapes, tables, charts) from a PPTX and "
            "store it with precise location metadata."
        )
    )
    parser.add_argument(
        "pptx",
        help="Path to the PPTX file to analyze.",
    )
    parser.add_argument(
        "--output",
        help="Destination JSON path; defaults to <pptx>.text-map.json",
    )
    return parser.parse_args()


def shape_type_name(shape) -> str:
    try:
        return MSO_SHAPE_TYPE(shape.shape_type).name
    except ValueError:
        return f"UNKNOWN_{shape.shape_type}"


def record_runs(
    text_frame,
    base_entry: Dict[str, Any],
    path_segments: Sequence[str],
    container: str,
    extra_fields: Dict[str, Any] | None = None,
    *,
    entries: List[Dict[str, Any]],
) -> None:
    extra_fields = extra_fields or {}
    tf_segments = [*path_segments, container]
    for para_idx, para in enumerate(text_frame.paragraphs, start=1):
        para_segments = [*tf_segments, f"paragraph[{para_idx}]"]
        for run_idx, run in enumerate(para.runs, start=1):
            text = run.text
            if not text:
                continue
            run_segments = [*para_segments, f"run[{run_idx}]"]
            entry = dict(base_entry)
            entry.update(extra_fields)
            entry.update(
                {
                    "container": container,
                    "paragraph_index": para_idx,
                    "run_index": run_idx,
                    "path": ".".join(run_segments),
                    "text": text,
                }
            )
            entries.append(entry)


def record_chart_text(
    shape, base_entry: Dict[str, Any], *, entries: List[Dict[str, Any]]
) -> None:
    if not getattr(shape, "has_chart", False):
        return
    chart_part = shape.chart.part
    nodes = chart_part._element.xpath(".//a:t")
    for chart_idx, node in enumerate(nodes, start=1):
        text = node.text
        if not text:
            continue
        entry = dict(base_entry)
        entry.update(
            {
                "container": "chart_part",
                "chart_partname": chart_part.partname,
                "chart_text_index": chart_idx,
                "chart_xpath": node.getroottree().getpath(node),
                "text": text,
            }
        )
        entries.append(entry)


def walk_shapes(
    shape_iter: Iterable[Any],
    path_segments: Sequence[str],
    index_chain: Sequence[int],
    slide_meta: Dict[str, Any],
    *,
    entries: List[Dict[str, Any]],
) -> None:
    for shape_idx, shape in enumerate(shape_iter, start=1):
        shape_segments = [*path_segments, f"shape[{shape_idx}]"]
        chain = [*index_chain, shape_idx]
        base_entry = {
            "slide_index": slide_meta["slide_index"],
            "slide_id": slide_meta["slide_id"],
            "slide_layout": slide_meta["slide_layout"],
            "shape_id": shape.shape_id,
            "shape_name": shape.name,
            "shape_type": shape_type_name(shape),
            "shape_index_chain": chain,
            "shape_path": ".".join(shape_segments),
        }

        if getattr(shape, "has_text_frame", False):
            record_runs(
                shape.text_frame,
                base_entry,
                shape_segments,
                "text_frame",
                entries=entries,
            )

        if getattr(shape, "has_table", False):
            table_segments = [*shape_segments, "table"]
            for row_idx, row in enumerate(shape.table.rows, start=1):
                for col_idx, cell in enumerate(row.cells, start=1):
                    cell_segments = [*table_segments, f"cell[{row_idx},{col_idx}]"]
                    extra = {
                        "table_row": row_idx,
                        "table_col": col_idx,
                    }
                    record_runs(
                        cell.text_frame,
                        base_entry,
                        cell_segments,
                        "table_cell",
                        extra,
                        entries=entries,
                    )

        record_chart_text(shape, base_entry, entries=entries)

        if shape.shape_type == MSO_SHAPE_TYPE.GROUP:
            walk_shapes(
                shape.shapes,
                shape_segments,
                chain,
                slide_meta,
                entries=entries,
            )


def build_text_map(pptx_path: Path) -> Dict[str, Any]:
    presentation = Presentation(pptx_path)
    entries: List[Dict[str, Any]] = []
    for slide_idx, slide in enumerate(presentation.slides, start=1):
        slide_meta = {
            "slide_index": slide_idx,
            "slide_id": slide.slide_id,
            "slide_layout": getattr(slide.slide_layout, "name", None),
        }
        walk_shapes(
            slide.shapes,
            [f"slide[{slide_idx}]"],
            [],
            slide_meta,
            entries=entries,
        )
    return {
        "source": str(pptx_path),
        "slide_count": len(presentation.slides),
        "entry_count": len(entries),
        "entries": entries,
    }


def main() -> None:
    args = parse_args()
    pptx_path = Path(args.pptx)
    assert pptx_path.exists(), f"PPTX path does not exist: {pptx_path}"
    output_path = (
        Path(args.output) if args.output else pptx_path.with_suffix(".text-map.json")
    )
    text_map = build_text_map(pptx_path)
    output_path.write_text(json.dumps(text_map, ensure_ascii=True, indent=2))


if __name__ == "__main__":
    main()
